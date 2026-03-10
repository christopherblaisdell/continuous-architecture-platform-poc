#!/usr/bin/env python3
"""
Generate full 12-page executive PowerPoint presentation for Continuous Architecture Platform.
Uses direct "our/we" company language for internal audiences.
Emphasizes vector database indexing as the key cost advantage.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Corporate color palette (teal/orange from the web theme)
TEAL = RGBColor(0, 137, 123)
ORANGE = RGBColor(255, 143, 0)
DARK_GRAY = RGBColor(55, 71, 79)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # All 12 slides
    add_title_slide(prs)                    # 1
    add_problem_slide(prs)                  # 2
    add_solution_overview_slide(prs)        # 3
    add_shared_workspace_slide(prs)         # 4
    add_markdown_first_slide(prs)           # 5
    add_cost_evidence_slide(prs)            # 6 - EMPHASIZE VECTOR DATABASE
    add_output_analysis_slide(prs)          # 7
    add_publishing_pipeline_slide(prs)      # 8
    add_live_demo_slide(prs)                # 9
    add_closing_loop_slide(prs)             # 10
    add_roadmap_slide(prs)                  # 11
    add_ask_slide(prs)                      # 12
    
    return prs

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Continuous Architecture Platform"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Building on Our Git-First Foundation"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Date and presenter placeholder
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "March 2026 • Architecture Practice"
    p = date_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Opening talking points:
- Full 12-page presentation covering all aspects of the platform
- This builds directly on our existing strengths - Git-controlled specs and production gating
- Everything I'll show you is already deployed and working
- Key innovation: Vector database indexing makes this 208x cheaper than alternatives
- Total time to present: 30-40 minutes with live demo
- Live portal demo included"""

def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "The Problem: The Last Mile"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # We already did the hard part
    hard_part_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    frame = hard_part_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "We Already Did the Hard Part:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    items = [
        "OpenAPI specs are source-controlled in Git",
        "PlantUML diagrams version-controlled alongside specs",
        "No Swagger goes to production without Git approval",
        "Production gating works - this is our foundation"
    ]
    
    for item in items:
        p = frame.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(18)
        p.level = 0
        p.space_before = Pt(6)
    
    # Two gaps
    gaps_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(3.5))
    frame = gaps_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "But Two Critical Gaps Erode This Foundation:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p = frame.add_paragraph()
    p.text = "Gap 1: Confluence Is Manual and Voluntary"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    p.font.color.rgb = DARK_GRAY
    
    p = frame.add_paragraph()
    p.text = "Specs are in Git. Confluence pages fall behind because updating them is optional."
    p.font.size = Pt(16)
    p.level = 1
    
    p = frame.add_paragraph()
    p.text = "Gap 2: Design vs Reality - No Reconciliation"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    p.font.color.rgb = DARK_GRAY
    
    p = frame.add_paragraph()
    p.text = "Developers deviate during implementation. Nobody verifies what was actually built matches what was designed."
    p.font.size = Pt(16)
    p.level = 1
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Key talking points:
- Emphasize that our Git-based foundation is already strong - most companies don't have this
- Gap 1: Confluence updates are voluntary, so they get skipped. The browsable docs stakeholders use fall behind.
- Gap 2: After deployment, we never verify the code matches the design. Specs describe intent, not reality.
- These gaps compound over time - every project makes it worse
- The next architect has to re-investigate because documentation is stale
- This presentation shows how we close both gaps"""

def add_solution_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "The Solution: Four Pillars"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Four pillars in 2x2 grid
    pillars = [
        ("1. Enhanced Workspace", 
         "Our Git repo + AI context\n500+ lines of domain knowledge\nPrevious designs inform new ones\nWorkspace grows richer over time"),
        ("2. Markdown Solution Designs",
         "Extends our Git-first practice\nADRs, impacts, guidance in Markdown\nPull request reviews on decisions\nAI-readable for future sessions"),
        ("3. Automated Publishing",
         "Replaces manual Confluence step\nGit push → portal updated\n301 artifacts auto-generated\nNo voluntary updates to skip"),
        ("4. AI-Assisted Workflows",
         "$39/month per seat (fixed cost)\nGitHub Copilot reads our Git repo\nProduces MADR ADRs, arc42 designs\n208x cheaper than alternatives")
    ]
    
    x_positions = [0.5, 5.0]
    y_positions = [1.2, 4.2]
    
    for i, (title, content) in enumerate(pillars):
        x = x_positions[i % 2]
        y = y_positions[i // 2]
        
        # Box background
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = TEAL
        box.line.width = Pt(2)
        
        # Title
        text_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(3.8), Inches(0.5))
        frame = text_box.text_frame
        frame.text = title
        p = frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEAL
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(3.8), Inches(1.5))
        frame = content_box.text_frame
        frame.word_wrap = True
        frame.text = content
        p = frame.paragraphs[0]
        p.font.size = Pt(14)
        p.line_spacing = 1.2
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Solution overview:
- All four pillars work together - each builds on our existing Git foundation
- Pillar 1 (Enhanced Workspace): Our Git repo becomes the AI's context - it reads our actual specs
- Pillar 2 (Markdown): We extend what we already do (Git-first) to include solution designs and ADRs
- Pillar 3 (Publishing): Automated publishing - a git push replaces the manual Confluence step
- Pillar 4 (AI Workflows): GitHub Copilot Pro+ with Claude Opus 4.6 - fixed subscription
- Next slides will dive deep into each pillar"""

def add_shared_workspace_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Pillar 1: Enhanced Workspace"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Key insight box
    insight_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(1.2))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = TEAL
    insight_box.line.width = Pt(0)
    
    insight_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.8))
    frame = insight_text.text_frame
    frame.word_wrap = True
    frame.text = "The AI doesn't hallucinate when it can read the actual source material. Our Git repo becomes the AI's shared context - same specs, diagrams, and source code we already maintain."
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # What's in the workspace
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(4.2))
    frame = content_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "What Lives in Our Architecture Repository:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    workspace_items = [
        ("19 OpenAPI specs", "Authoritative API contracts for every microservice"),
        ("11 Architecture Decision Records", "Global decision log with rationale"),
        ("6 Service architecture pages", "Living baselines per service"),
        ("8 Java source files", "Key service implementations for analysis"),
        ("3 Mock tool scripts", "JIRA, Elasticsearch, GitLab simulation"),
        ("4 Architecture standards", "arc42, MADR, C4 model, ISO 25010"),
        ("copilot-instructions.md", "500+ lines of domain knowledge loaded into every AI session")
    ]
    
    for item, description in workspace_items:
        p = frame.add_paragraph()
        p.text = f"{item}: {description}"
        p.font.size = Pt(14)
        p.space_before = Pt(8)
        p.level = 0
        # Bold the item name
        if p.runs:
            p.runs[0].font.bold = True
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Enhanced Workspace deep-dive:
- This is the foundation - our Git repo becomes the AI's context
- AI reads the ACTUAL specs we already maintain, not synthetic examples
- copilot-instructions.md is the secret weapon - 500+ lines of domain knowledge auto-loaded
- Includes: role definition, domain model, bounded context rules, anti-patterns, standards
- The workspace gets richer with every project - and Copilot's fixed pricing means more context costs nothing
- This is why accuracy is high - AI has full context, doesn't need to guess"""

def add_markdown_first_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Pillar 2: Markdown Solution Designs"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Extending what we do
    extend_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.8))
    frame = extend_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Extending Our Git-First Practice:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    p = frame.add_paragraph()
    p.text = "We already author OpenAPI specs in YAML and diagrams in PlantUML - text formats checked into Git. We're extending this to solution designs, ADRs, and impact assessments in Markdown."
    p.font.size = Pt(16)
    p.space_before = Pt(6)
    
    # What we have vs what we add - table
    table_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(2.8))
    frame = table_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "What We Have vs What We Add:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    items = [
        ("OpenAPI specs", "Already in Git (gated)", "→ Same + AI-assisted updates"),
        ("PlantUML diagrams", "Already in Git", "→ Same + auto-rendered to portal"),
        ("Solution designs", "Confluence (manual)", "→ Markdown in Git + auto-published"),
        ("Architecture Decision Records", "Ticket branches (not discoverable)", "→ Markdown in Git + global log"),
        ("Impact assessments", "Confluence or email", "→ Markdown in Git + version-controlled"),
        ("Service documentation", "Confluence (voluntary)", "→ Auto-generated from specs")
    ]
    
    for artifact, today, future in items:
        p = frame.add_paragraph()
        p.text = f"{artifact}: {today} {future}"
        p.font.size = Pt(13)
        p.space_before = Pt(6)
    
    # Benefits
    benefits_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    frame = benefits_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Benefits:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    benefits = [
        "Pull request reviews on architecture decisions, not just specs",
        "AI-readable history - the model analyzes past designs automatically",
        "Searchable decision log - ADRs indexed and discoverable",
        "Publishable everywhere - MkDocs portal, Confluence API sync, PDF export from one source"
    ]
    
    for benefit in benefits:
        p = frame.add_paragraph()
        p.text = "• " + benefit
        p.font.size = Pt(14)
        p.space_before = Pt(4)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Markdown-first authoring:
- This extends what we already do - we're already Git-first for specs and diagrams
- Now adding: solution designs, ADRs, impact assessments, implementation guidance
- Everything in Markdown, everything version-controlled
- Pull request workflow: architects review designs with line-by-line diffs, just like code
- AI benefit: Markdown is AI-readable. Past designs become context for future sessions.
- No more buried ADRs in ticket branches - they get promoted to a global searchable log
- Single source of truth: one Markdown file publishes to portal, Confluence, and PDF"""

def add_cost_evidence_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Cost Evidence: The Vector Database Advantage"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # HUGE CALLOUT BOX - Vector Database
    vector_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(2.2))
    vector_box.fill.solid()
    vector_box.fill.fore_color.rgb = ORANGE
    vector_box.line.width = Pt(3)
    vector_box.line.color.rgb = DARK_GRAY
    
    vector_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.5))
    frame = vector_title.text_frame
    frame.text = "THE ARCHITECTURAL DIFFERENCE"
    p = frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    vector_content = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(8.6), Inches(1.2))
    frame = vector_content.text_frame
    frame.word_wrap = True
    frame.text = "GitHub Copilot pre-indexes our entire workspace into a VECTOR DATABASE. All our specs, diagrams, previous designs, domain knowledge - indexed once, used unlimited times. This is amortized across our fixed $39/month subscription.\n\nOpenRouter recalculates everything from scratch on every request. Pays per token, every time."
    p = frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.3
    
    # Cost comparison - side by side
    # GitHub Copilot side
    copilot_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(4.3), Inches(2.8))
    frame = copilot_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "GitHub Copilot Pro+"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    p = frame.add_paragraph()
    p.text = "$39/month fixed subscription"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_before = Pt(8)
    
    copilot_points = [
        "Vector database indexes workspace once",
        "Richer context costs nothing extra",
        "Intent-based billing: only human prompts count",
        "Autonomous tool calls are FREE",
        "As workspace grows, cost stays flat"
    ]
    
    for point in copilot_points:
        p = frame.add_paragraph()
        p.text = "✓ " + point
        p.font.size = Pt(13)
        p.space_before = Pt(4)
    
    # OpenRouter side
    openrouter_box = slide.shapes.add_textbox(Inches(5.2), Inches(3.7), Inches(4.3), Inches(2.8))
    frame = openrouter_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "OpenRouter (Per-Token)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    
    p = frame.add_paragraph()
    p.text = "$155.73 per run (variable)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_before = Pt(8)
    
    openrouter_points = [
        "Recalculates context every request",
        "Pays for every token, every time",
        "No indexing - starts from scratch",
        "Context size = direct cost",
        "As workspace grows, cost scales linearly"
    ]
    
    for point in openrouter_points:
        p = frame.add_paragraph()
        p.text = "✗ " + point
        p.font.size = Pt(13)
        p.space_before = Pt(4)
    
    # Bottom: 208x
    result_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.6))
    frame = result_box.text_frame
    frame.text = "Result: 208x cost advantage per run, and the gap widens as our workspace grows"
    p = frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """CRITICAL TALKING POINT - Vector Database Advantage:
- This is THE most important slide - explains why 208x difference exists
- GitHub Copilot: Pre-indexes our workspace into a vector database (semantic search index)
- Indexing happens once, used unlimited times - amortized across fixed subscription
- OpenRouter: No indexing. Recalculates entire context from scratch every request, pays per token
- Intent-based billing: In Copilot, only human-typed prompts count. All autonomous tool calls are free.
- As our workspace grows (more ADRs, more designs), Copilot cost stays $39/month flat
- OpenRouter cost scales linearly with workspace size - bigger context = higher per-run cost
- This is a structural architectural advantage, not a temporary pricing difference
- Vector database is how Copilot achieves both low cost AND high accuracy"""

def add_output_analysis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Output Analysis: What the AI Produced"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Big number
    number_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(1.5))
    frame = number_box.text_frame
    frame.text = "39"
    p = frame.paragraphs[0]
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER
    
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4), Inches(0.4))
    frame = label_box.text_frame
    frame.text = "complete architecture files"
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # Details
    details_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(5.5))
    frame = details_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "From 5 Complex Scenarios:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    outputs = [
        "Solution designs following arc42 structure",
        "Impact assessments for all affected services",
        "MADR-formatted Architecture Decision Records",
        "Implementation guidance with code patterns",
        "User stories with acceptance criteria",
        "PlantUML C4 diagrams with proper notation",
        "OpenAPI spec updates from approved designs",
        "Investigation reports grounded in logs and source code"
    ]
    
    for output in outputs:
        p = frame.add_paragraph()
        p.text = "✓ " + output
        p.font.size = Pt(14)
        p.space_before = Pt(8)
    
    # Accuracy note
    accuracy_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(4), Inches(3))
    frame = accuracy_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Accuracy Assessment:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p = frame.add_paragraph()
    p.text = "Zero fabricated data - all analysis grounded in workspace evidence"
    p.font.size = Pt(15)
    p.space_before = Pt(10)
    
    p = frame.add_paragraph()
    p.text = "MADR, arc42, and C4 model standards followed correctly"
    p.font.size = Pt(15)
    p.space_before = Pt(8)
    
    p = frame.add_paragraph()
    p.text = "Completeness: 39 files vs 37 from per-token alternative"
    p.font.size = Pt(15)
    p.space_before = Pt(8)
    
    p = frame.add_paragraph()
    p.text = "Accuracy: 0 hallucinated fields vs 4 from per-token alternative"
    p.font.size = Pt(15)
    p.space_before = Pt(8)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Output quality talking points:
- 39 files across 5 scenarios - investigation, design, implementation guidance, decisions
- Every artifact follows our documented standards (MADR for ADRs, arc42 for solution designs)
- No fabricated data - AI only used evidence from logs, specs, and source code
- Compared to per-token alternative: Copilot produced 39 files vs 37, and had zero fabrication vs 4 hallucinated fields
- These aren't toy examples - complex multi-service scenarios like RFID wristband integration and schedule overwrite bugs
- Full architectural depth - diagrams, trade-off analysis, consequences documented
- The vector database indexing gives AI full context, so it doesn't need to hallucinate"""

def add_publishing_pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Pillar 3: Automated Publishing Pipeline"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Key message
    message_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.7))
    frame = message_box.text_frame
    frame.word_wrap = True
    frame.text = "Replaces the voluntary manual Confluence step with: git push → portal automatically updated"
    p = frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    
    # What gets published
    published_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.5), Inches(4.5))
    frame = published_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "What Gets Auto-Published:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    items = [
        "19 microservice deep-dive pages",
        "139 endpoint sequence diagrams (SVG)",
        "20 C4 context diagrams",
        "19 Swagger UI pages (interactive)",
        "11 Architecture Decision Records",
        "Full-text search across all docs",
        "Cross-service navigation with deep linking"
    ]
    
    for item in items:
        p = frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.space_before = Pt(8)
    
    # Before/After
    comparison_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.3), Inches(4.5))
    frame = comparison_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Before vs After:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    p = frame.add_paragraph()
    p.text = "BEFORE:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(12)
    p.font.color.rgb = RGBColor(150, 0, 0)
    
    p = frame.add_paragraph()
    p.text = "Voluntary Confluence update (often skipped)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = frame.add_paragraph()
    p.text = "Manual diagram re-rendering"
    p.font.size = Pt(14)
    p.level = 1
    
    p = frame.add_paragraph()
    p.text = "Broken cross-service links"
    p.font.size = Pt(14)
    p.level = 1
    
    p = frame.add_paragraph()
    p.text = "AFTER:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(16)
    p.font.color.rgb = TEAL
    
    p = frame.add_paragraph()
    p.text = "git push (automatic, never skipped)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = frame.add_paragraph()
    p.text = "Auto-rendered diagrams"
    p.font.size = Pt(14)
    p.level = 1
    
    p = frame.add_paragraph()
    p.text = "Auto-generated navigation"
    p.font.size = Pt(14)
    p.level = 1
    
    # Cost
    cost_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.4))
    frame = cost_box.text_frame
    frame.text = "Publishing Cost: $0 (Azure Static Web Apps free tier + MkDocs Material open source)"
    p = frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Automated publishing pipeline:
- This closes Gap 1 (Confluence falls behind)
- Current problem: updating Confluence is manual and voluntary, so it gets skipped
- Solution: git push triggers automated build and deployment
- Build steps: Python generators read specs, produce pages, PlantUML renders diagrams, MkDocs builds site
- Build time: under 30 seconds from push to live
- 301 total artifacts published automatically - microservice pages, diagrams, Swagger UI, ADR log
- Optional: can sync to Confluence via API if needed - same Markdown source, two renderers
- Cost: $0 (free tier Azure + open source tools)
- Result: browsable documentation always matches Git state, never falls behind"""

def add_live_demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Live Demo: The Architecture Portal"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Big CTA box
    cta_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(2))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = TEAL
    cta_box.line.width = Pt(0)
    
    # URL text
    url_box = slide.shapes.add_textbox(Inches(1.7), Inches(1.8), Inches(6.6), Inches(1.4))
    frame = url_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Live Demo:"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = frame.add_paragraph()
    p.text = "https://architecture.novatrek.cc"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    
    # Make it a hyperlink
    run = p.runs[0]
    run.hyperlink.address = "https://architecture.novatrek.cc"
    
    # What to explore
    explore_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(3))
    frame = explore_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "What to Explore in the Live Portal:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    features = [
        "19 microservice deep-dive pages - auto-generated from OpenAPI specs",
        "139 clickable SVG sequence diagrams - every endpoint documented with deep linking",
        "Interactive Swagger UI - try-it-out capable API documentation",
        "Enterprise C4 diagram - all services with PCI compliance zone highlighted",
        "Global architecture decision log - 11 ADRs with cross-references",
        "Full-text search - instant results across all documentation"
    ]
    
    for feature in features:
        p = frame.add_paragraph()
        p.text = "• " + feature
        p.font.size = Pt(16)
        p.space_before = Pt(8)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Portal demonstration talking points:
- PAUSE HERE and open the portal in a browser
- This is a LIVE site, deployed to Azure Static Web Apps
- Everything you'll see was generated by AI from our OpenAPI specs
- Walk through: 
  1. Home page - service catalog
  2. Click svc-check-in (most connected service)
  3. Show sequence diagram for an endpoint - it's clickable SVG
  4. Click a cross-service integration arrow - deep links to target endpoint
  5. Show Swagger UI - interactive, try-it-out capable
  6. Show global ADR log - 11 decisions with full context
- Emphasize: this publishes automatically on git push - no manual Confluence step
- All 301 artifacts from a single git push"""

def add_closing_loop_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Pillar 4: Closing the Loop (PROMOTE)"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # The problem we've never solved
    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(1))
    frame = problem_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Gap 2: The Problem We've Never Solved"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p = frame.add_paragraph()
    p.text = "After deployment, we never reconcile what was BUILT vs what was DESIGNED. Specs describe intent. Production code may differ. Nobody closes this gap."
    p.font.size = Pt(16)
    p.space_before = Pt(8)
    
    # The workflow
    workflow_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.8))
    frame = workflow_box.text_frame
    frame.text = "INTAKE → INVESTIGATE → DESIGN → BUILD → DEPLOY → PROMOTE → DONE"
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER
    
    # What PROMOTE does
    promote_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(3.2))
    frame = promote_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "What the PROMOTE Step Does:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    actions = [
        "Reconcile specs against reality - verify what was actually built",
        "Update OpenAPI specs to reflect production behavior, not just design intent",
        "Promote ticket-level ADRs to the global searchable decision log",
        "Update service baseline pages with current state",
        "Create audit trail of design → reality reconciliation",
        "Make architectural knowledge discoverable for the next project"
    ]
    
    for action in actions:
        p = frame.add_paragraph()
        p.text = "• " + action
        p.font.size = Pt(15)
        p.space_before = Pt(8)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """PROMOTE step innovation:
- This closes Gap 2 (design vs reality divergence)
- This is the architectural innovation that makes it "continuous"
- Current reality: architects create designs, developers implement, nobody verifies alignment
- After every project: ADRs stay in ticket branches, specs describe intent not reality, knowledge is lost
- PROMOTE step: AI reconciles design vs implementation, updates specs to match reality, promotes ADRs to global log
- This step has never existed in our workflow before - no time, no process, no tooling
- Now automated: AI does the reconciliation work, creates audit trail
- Result: specs describe actual production behavior, not just original design intent
- The vector database makes this feasible - AI can read specs, source code, and logs to compare"""

def add_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Roadmap: Six Phases"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Phases
    phases_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(5.8))
    frame = phases_box.text_frame
    frame.word_wrap = True
    
    phases = [
        ("Phase 1: AI Toolchain Comparison", "COMPLETE - March 2026", TEAL, 
         "Cost comparison, output analysis, toolchain recommendation (Copilot Pro+)"),
        ("Phase 2: AI-Integrated Workflow Design", "April - May 2026", DARK_GRAY,
         "Standardized AI workflows, PROMOTE playbook, execution prompt library"),
        ("Phase 3: Pipeline Integration", "May - June 2026", DARK_GRAY,
         "CI/CD validation gates, automated spec linting, quality checks"),
        ("Phase 4: Navigable Artifact Graph", "June - July 2026", DARK_GRAY,
         "Bidirectional traceability, artifact graph, backlink generation"),
        ("Phase 5: Continuous Improvement Loop", "July - August 2026", DARK_GRAY,
         "Quality metrics, pattern identification, AI instruction optimization"),
        ("Phase 6: Documentation Publishing", "COMPLETE - March 2026", TEAL,
         "Live architecture portal with automated publishing pipeline")
    ]
    
    y_offset = 0
    for title, timeline, color, description in phases:
        p = frame.add_paragraph()
        p.text = f"{title} - {timeline}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        if y_offset > 0:
            p.space_before = Pt(12)
        
        p = frame.add_paragraph()
        p.text = description
        p.font.size = Pt(13)
        p.level = 1
        p.space_before = Pt(4)
        
        y_offset += 1
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Roadmap overview:
- Six phases total, each delivers independent value
- Phase 1 (COMPLETE): Compared GitHub Copilot vs Roo Code, established 208x cost advantage
- Phase 6 (COMPLETE): Built the live portal you just saw, automated publishing pipeline
- We accelerated Phase 6 to demonstrate the full platform capability
- Phase 2 (Next): Formalize workflows - make PROMOTE a standard step, create playbooks
- Phase 3: Integrate into CI/CD - automated validation on every commit
- Phase 4: Full traceability - from any ADR, navigate to ticket, design, impacted specs
- Phase 5: Continuous improvement - measure quality, optimize AI instructions
- Timeline: Full platform by August 2026
- Can stop at any phase and still have value"""

def add_ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "The Ask: What We Need to Proceed"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Four approvals in grid
    approvals = [
        ("1. GitHub Copilot Licenses",
         "$39/month per seat\nFixed subscription - no variable costs\nVector database indexing\nStart with architecture team pilot"),
        ("2. Markdown in Git",
         "Extend our Git-first practice\nSolution designs in Markdown\nPull request reviews on ADRs\nAI reads designs in future sessions"),
        ("3. Automated Publishing",
         "Azure Static Web Apps (free tier)\ngit push → portal updated\nReplaces manual Confluence step\n$0 infrastructure cost"),
        ("4. Optional: Confluence API",
         "Sync Markdown to Confluence\nGit remains source of truth\nConfluence becomes read-only mirror\nNo dual maintenance required")
    ]
    
    x_positions = [0.5, 5.0]
    y_positions = [1.2, 3.8]
    
    for i, (title, content) in enumerate(approvals):
        x = x_positions[i % 2]
        y = y_positions[i // 2]
        
        # Box
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = TEAL
        box.line.width = Pt(2)
        
        # Title
        text_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(3.8), Inches(0.4))
        frame = text_box.text_frame
        frame.text = title
        p = frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEAL
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.6), Inches(3.8), Inches(1.3))
        frame = content_box.text_frame
        frame.word_wrap = True
        frame.text = content
        p = frame.paragraphs[0]
        p.font.size = Pt(13)
        p.line_spacing = 1.2
    
    # Total cost and next steps
    bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(1))
    frame = bottom_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Total Monthly Cost: $39/seat (Copilot) + $0 (Azure free tier) = $39/seat"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    p = frame.add_paragraph()
    p.text = "Next Steps: Pilot with architecture team (1 month) → Measure quality → Expand if successful"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Closing and approval request:
- We need four approvals to adopt this platform
- Approval 1: Copilot licenses - THE KEY INVESTMENT. Vector database indexing is the secret sauce.
  * Start with architecture team as pilot
  * Measure quality and adoption for 1 month
  * Expand if successful
- Approval 2: Markdown authoring - extends what we already do (Git-first)
  * Just adds solution designs to the repo alongside specs
- Approval 3: Automated publishing - replaces the manual Confluence step that gets skipped anyway
  * $0 cost (free tier Azure + open source MkDocs)
- Approval 4: Optional - Confluence sync if we need to maintain Confluence as canonical platform
  * Git remains source of truth, Confluence becomes read-only mirror
- Total cost: $39/month per architect (fixed) - no variable costs, no infrastructure expenses
- Everything shown today is already built and deployed - this isn't vaporware
- Questions to address: pilot scope, timeline, success metrics, governance
- Emphasize: The vector database is what makes this work - low cost AND high accuracy"""

def save_presentation(prs, output_path):
    prs.save(output_path)
    print(f"PowerPoint created: {output_path}")

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "../presentation-internal-full.pptx"
    save_presentation(prs, output_path)
    print("\nPresentation Details:")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Format: Full presentation (all 12 pages)")
    print(f"  Language: Internal (our/we)")
    print(f"  Speaker notes: Included on every slide")
    print(f"  Live portal link: Included on slide 9")
    print(f"  Vector database: PROMINENTLY featured on slide 6")
