#!/usr/bin/env python3
"""
Generate executive PowerPoint presentation for Continuous Architecture Platform.
Uses direct "our/we" company language for internal audiences.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Corporate color palette (teal/orange from the web theme)
TEAL = RGBColor(0, 137, 123)
ORANGE = RGBColor(255, 143, 0)
DARK_GRAY = RGBColor(55, 71, 79)
WHITE = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: The Problem
    add_problem_slide(prs)
    
    # Slide 3: The Solution
    add_solution_slide(prs)
    
    # Slide 4: Cost Evidence
    add_cost_slide(prs)
    
    # Slide 5: What We Built
    add_demo_slide(prs)
    
    # Slide 6: Output Analysis
    add_output_slide(prs)
    
    # Slide 7: The PROMOTE Innovation
    add_promote_slide(prs)
    
    # Slide 8: The Ask
    add_ask_slide(prs)
    
    return prs

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Continuous Architecture Platform"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Building on Our Git-First Foundation"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Date and presenter placeholder
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "March 2026 • Architecture Practice"
    p = date_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Opening talking points:
- We've built a proof of concept that transforms how our architecture practice works
- This builds directly on our existing strengths - Git-controlled specs and production gating
- Everything I'll show you is already deployed and working
- Total time to present: 15-20 minutes
- Live portal demo included"""

def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "The Problem: The Last Mile"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # We already did the hard part
    hard_part_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    frame = hard_part_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "We Already Did the Hard Part:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    items = [
        "OpenAPI specs are source-controlled in Git",
        "PlantUML diagrams version-controlled alongside specs",
        "No Swagger goes to production without Git approval",
        "Production gating works - this is our foundation"
    ]
    
    for item in items:
        p = frame.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(18)
        p.level = 0
        p.space_before = Pt(6)
    
    # Two gaps
    gaps_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(3.5))
    frame = gaps_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "But Two Critical Gaps Erode This Foundation:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p = frame.add_paragraph()
    p.text = "Gap 1: Confluence Is Manual and Voluntary"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    p.font.color.rgb = DARK_GRAY
    
    p = frame.add_paragraph()
    p.text = "Specs are in Git. Confluence pages fall behind because updating them is optional."
    p.font.size = Pt(16)
    p.level = 1
    
    p = frame.add_paragraph()
    p.text = "Gap 2: Design vs Reality - No Reconciliation"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    p.font.color.rgb = DARK_GRAY
    
    p = frame.add_paragraph()
    p.text = "Developers deviate during implementation. Nobody verifies what was actually built matches what was designed."
    p.font.size = Pt(16)
    p.level = 1
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Key talking points:
- Emphasize that our Git-based foundation is already strong - most companies don't have this
- Gap 1: Confluence updates are voluntary, so they get skipped. The browsable docs stakeholders use fall behind.
- Gap 2: After deployment, we never verify the code matches the design. Specs describe intent, not reality.
- These gaps compound over time - every project makes it worse
- The next architect has to re-investigate because documentation is stale"""

def add_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "The Solution: Four Pillars"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Four pillars in 2x2 grid
    pillars = [
        ("1. AI-Assisted Workflows", 
         "$39/month per seat\nGitHub Copilot reads our Git repo\nProduces MADR ADRs, arc42 designs\n208x cheaper than alternatives"),
        ("2. Enhanced Workspace",
         "Our Git repo + AI context\n500+ lines of domain knowledge\nPrevious designs inform new ones\nWorkspace grows richer over time"),
        ("3. Markdown Solution Designs",
         "Extends our Git-first practice\nADRs, impacts, guidance in Markdown\nPull request reviews on decisions\nAI-readable for future sessions"),
        ("4. Automated Publishing",
         "Replaces manual Confluence step\nGit push → portal updated\n301 artifacts auto-generated\nNo voluntary updates to skip")
    ]
    
    x_positions = [0.5, 5.0]
    y_positions = [1.2, 4.2]
    
    for i, (title, content) in enumerate(pillars):
        x = x_positions[i % 2]
        y = y_positions[i // 2]
        
        # Box background
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = TEAL
        box.line.width = Pt(2)
        
        # Title
        text_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(3.8), Inches(0.5))
        frame = text_box.text_frame
        frame.text = title
        p = frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEAL
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(3.8), Inches(1.5))
        frame = content_box.text_frame
        frame.word_wrap = True
        frame.text = content
        p = frame.paragraphs[0]
        p.font.size = Pt(14)
        p.line_spacing = 1.2
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Solution overview:
- Pillar 1: We're using GitHub Copilot Pro+ with Claude Opus 4.6 - fixed subscription, no variable costs
- Pillar 2: Our Git repo becomes the AI's context - it reads our actual specs, not synthetic examples
- Pillar 3: We extend what we already do (Git-first) to include solution designs and ADRs in Markdown
- Pillar 4: Automated publishing - a git push replaces the manual Confluence step that gets skipped
- All four pillars work together - each builds on our existing Git foundation"""

def add_cost_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Cost Evidence: The Vector Database Advantage"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # HUGE CALLOUT BOX - Vector Database
    vector_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(2.2))
    vector_box.fill.solid()
    vector_box.fill.fore_color.rgb = ORANGE
    vector_box.line.width = Pt(3)
    vector_box.line.color.rgb = DARK_GRAY
    
    vector_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.5))
    frame = vector_title.text_frame
    frame.text = "THE ARCHITECTURAL DIFFERENCE"
    p = frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    vector_content = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(8.6), Inches(1.2))
    frame = vector_content.text_frame
    frame.word_wrap = True
    frame.text = "GitHub Copilot pre-indexes our entire workspace into a VECTOR DATABASE. All our specs, diagrams, previous designs, domain knowledge - indexed once, used unlimited times. This is amortized across our fixed $39/month subscription.\n\nOpenRouter recalculates everything from scratch on every request. Pays per token, every time."
    p = frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.3
    
    # Cost comparison - side by side
    # GitHub Copilot side
    copilot_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(4.3), Inches(2.8))
    frame = copilot_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "GitHub Copilot Pro+"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    p = frame.add_paragraph()
    p.text = "$39/month fixed subscription"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_before = Pt(8)
    
    copilot_points = [
        "Vector database indexes workspace once",
        "Richer context costs nothing extra",
        "Intent-based billing: only human prompts count",
        "Autonomous tool calls are FREE",
        "As workspace grows, cost stays flat"
    ]
    
    for point in copilot_points:
        p = frame.add_paragraph()
        p.text = "✓ " + point
        p.font.size = Pt(13)
        p.space_before = Pt(4)
    
    # OpenRouter side
    openrouter_box = slide.shapes.add_textbox(Inches(5.2), Inches(3.7), Inches(4.3), Inches(2.8))
    frame = openrouter_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "OpenRouter (Per-Token)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    
    p = frame.add_paragraph()
    p.text = "$155.73 per run (variable)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_before = Pt(8)
    
    openrouter_points = [
        "Recalculates context every request",
        "Pays for every token, every time",
        "No indexing - starts from scratch",
        "Context size = direct cost",
        "As workspace grows, cost scales linearly"
    ]
    
    for point in openrouter_points:
        p = frame.add_paragraph()
        p.text = "✗ " + point
        p.font.size = Pt(13)
        p.space_before = Pt(4)
    
    # Bottom: 208x
    result_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.6))
    frame = result_box.text_frame
    frame.text = "Result: 208x cost advantage per run, and the gap widens as our workspace grows"
    p = frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """CRITICAL TALKING POINT - Vector Database Advantage:
- This is THE most important slide - explains why 208x difference exists
- GitHub Copilot: Pre-indexes our workspace into a vector database (semantic search index)
- Indexing happens once, used unlimited times - amortized across fixed subscription
- OpenRouter: No indexing. Recalculates entire context from scratch every request, pays per token
- Intent-based billing: In Copilot, only human-typed prompts count. All autonomous tool calls are free.
- As our workspace grows (more ADRs, more designs), Copilot cost stays $39/month flat
- OpenRouter cost scales linearly with workspace size - bigger context = higher per-run cost
- This is a structural architectural advantage, not a temporary pricing difference
- Vector database is how Copilot achieves both low cost AND high accuracy"""

def add_demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "What We Built: Live Portal"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Big CTA box
    cta_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(2))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = TEAL
    cta_box.line.width = Pt(0)
    
    # URL text
    url_box = slide.shapes.add_textbox(Inches(1.7), Inches(1.8), Inches(6.6), Inches(1.4))
    frame = url_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Live Demo:"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = frame.add_paragraph()
    p.text = "https://mango-sand-083b8ce0f.4.azurestaticapps.net"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    
    # Make it a hyperlink (will work in PowerPoint)
    run = p.runs[0]
    run.hyperlink.address = "https://mango-sand-083b8ce0f.4.azurestaticapps.net"
    
    # What to explore
    explore_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(3))
    frame = explore_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "What to Explore:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    features = [
        "19 microservice deep-dive pages - auto-generated from our OpenAPI specs",
        "139 clickable SVG sequence diagrams - every endpoint documented",
        "Interactive Swagger UI - try-it-out capable API documentation",
        "Enterprise C4 diagram - all services with PCI compliance zone highlighted",
        "Global architecture decision log - 11 ADRs with cross-references",
        "Full-text search - instant results across all documentation"
    ]
    
    for feature in features:
        p = frame.add_paragraph()
        p.text = "• " + feature
        p.font.size = Pt(16)
        p.space_before = Pt(8)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Portal demonstration talking points:
- This is a LIVE site, deployed to Azure Static Web Apps
- Everything you'll see was generated by AI from our specs
- Uses NovaTrek Adventures as synthetic domain (no real corporate data)
- Click through to show: microservice page with clickable diagrams, Swagger UI, global ADR log
- Emphasize: this publishes automatically on git push - no manual Confluence step
- Point out: deep linking works - diagrams link to specific endpoints on other service pages
- 301 total artifacts auto-published from a single git push"""

def add_output_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "What the AI Produced"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Big number
    number_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(1.5))
    frame = number_box.text_frame
    frame.text = "39"
    p = frame.paragraphs[0]
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER
    
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4), Inches(0.4))
    frame = label_box.text_frame
    frame.text = "complete architecture files"
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # Details
    details_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(5.5))
    frame = details_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "From 5 Complex Scenarios:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    outputs = [
        "Solution designs following arc42 structure",
        "Impact assessments for all affected services",
        "MADR-formatted Architecture Decision Records",
        "Implementation guidance with code patterns",
        "User stories with acceptance criteria",
        "PlantUML C4 diagrams with proper notation",
        "OpenAPI spec updates from approved designs",
        "Investigation reports grounded in logs and source code"
    ]
    
    for output in outputs:
        p = frame.add_paragraph()
        p.text = "✓ " + output
        p.font.size = Pt(14)
        p.space_before = Pt(8)
    
    # Accuracy note
    accuracy_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4), Inches(2))
    frame = accuracy_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "Accuracy Assessment:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p = frame.add_paragraph()
    p.text = "Zero fabricated data - all analysis grounded in workspace evidence"
    p.font.size = Pt(14)
    p.space_before = Pt(8)
    
    p = frame.add_paragraph()
    p.text = "MADR, arc42, and C4 model standards followed correctly"
    p.font.size = Pt(14)
    p.space_before = Pt(6)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Output quality talking points:
- 39 files across 5 scenarios - investigation, design, implementation guidance, decisions
- Every artifact follows our documented standards (MADR for ADRs, arc42 for solution designs)
- No fabricated data - AI only used evidence from logs, specs, and source code
- Compared to per-token alternative: Copilot produced 39 files vs 37, and had zero fabrication vs 4 hallucinated fields
- These aren't toy examples - complex multi-service scenarios like RFID wristband integration and schedule overwrite bugs
- Full architectural depth - diagrams, trade-off analysis, consequences documented"""

def add_promote_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "The Innovation: PROMOTE Step"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # The problem
    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    frame = problem_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "The Problem We've Never Solved:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p = frame.add_paragraph()
    p.text = "After deployment, we never reconcile what was built vs what was designed. Specs describe intent. Production code may differ. Nobody closes this gap."
    p.font.size = Pt(16)
    p.space_before = Pt(8)
    
    # The workflow
    workflow_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
    frame = workflow_box.text_frame
    frame.text = "INTAKE → INVESTIGATE → DESIGN → BUILD → DEPLOY → PROMOTE → DONE"
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER
    
    # What PROMOTE does
    promote_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.8))
    frame = promote_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "What PROMOTE Does:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    actions = [
        "Reconcile specs against reality - verify what was actually built",
        "Update OpenAPI specs to reflect production behavior, not just design intent",
        "Promote ticket-level ADRs to the global searchable decision log",
        "Update service baseline pages with current state",
        "Create audit trail of design → reality reconciliation"
    ]
    
    for action in actions:
        p = frame.add_paragraph()
        p.text = "• " + action
        p.font.size = Pt(16)
        p.space_before = Pt(8)
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """PROMOTE step innovation:
- This is the architectural innovation that makes it "continuous"
- Current reality: architects create designs, developers implement, nobody verifies alignment
- After every project: ADRs stay in ticket branches, specs describe intent not reality, knowledge is lost
- PROMOTE step: AI reconciles design vs implementation, updates specs to match reality, promotes ADRs to global log
- This step has never existed in our workflow before - no time, no process, no tooling
- Now automated: AI does the reconciliation work, creates audit trail
- Result: specs describe actual production behavior, not just original design intent"""

def add_ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "The Ask: What We Need"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Four approvals in grid
    approvals = [
        ("1. GitHub Copilot Licenses",
         "$39/month per seat\nFixed subscription\nNo variable costs\nStart with architecture team"),
        ("2. Markdown in Git",
         "Extend our Git-first practice\nSolution designs in Markdown\nPull request reviews on ADRs\nAI reads designs in future sessions"),
        ("3. Automated Publishing",
         "Azure Static Web Apps (free tier)\nGit push → portal updated\nReplaces manual Confluence step\nNo voluntary updates to skip"),
        ("4. Optional: Confluence API",
         "Sync Markdown to Confluence\nGit remains source of truth\nConfluence becomes read-only mirror\nNo dual maintenance")
    ]
    
    x_positions = [0.5, 5.0]
    y_positions = [1.2, 4.0]
    
    for i, (title, content) in enumerate(approvals):
        x = x_positions[i % 2]
        y = y_positions[i // 2]
        
        # Box
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = TEAL
        box.line.width = Pt(2)
        
        # Title
        text_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(3.8), Inches(0.4))
        frame = text_box.text_frame
        frame.text = title
        p = frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEAL
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.6), Inches(3.8), Inches(1.5))
        frame = content_box.text_frame
        frame.word_wrap = True
        frame.text = content
        p = frame.paragraphs[0]
        p.font.size = Pt(13)
        p.line_spacing = 1.2
    
    # Total cost
    cost_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.5))
    frame = cost_box.text_frame
    frame.text = "Total Monthly Cost: $39/seat (Copilot) + $0 (Azure free tier) + $0 (MkDocs)"
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """Closing and approval request:
- We need four approvals to adopt this platform
- Approval 1: Copilot licenses - start with architecture team as pilot, measure quality, expand if successful
- Approval 2: Markdown authoring - this extends what we already do (Git-first), just adds solution designs to the repo
- Approval 3: Automated publishing - replaces the manual Confluence step that gets skipped anyway
- Approval 4: Optional - Confluence sync if we need to maintain Confluence as canonical platform
- Total cost: $39/month per architect (fixed) - no infrastructure costs, no variable expenses
- Everything shown today is already built and deployed - this isn't vaporware
- Questions to address: pilot scope, timeline, success metrics, reporting structure"""

def save_presentation(prs, output_path):
    prs.save(output_path)
    print(f"PowerPoint created: {output_path}")

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "../presentation-internal.pptx"
    save_presentation(prs, output_path)
    print("\nPresentation Details:")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Format: Executive pitch (condensed)")
    print(f"  Language: Internal (our/we)")
    print(f"  Speaker notes: Included on every slide")
    print(f"  Live portal link: Included on slide 5")
